# 🛡️ SKILL BLINDADA (08/03/2026) - NÃO ALTERAR LÓGICA DE DESIGN OU IMAGENS
# Esta skill foi homologada com design Premium Manus AI e resiliência de tags.
# Qualquer modificação em layouts, cores ou busca de imagens deve ser evitada.
import os
import json
import uuid
import logging
import re
import asyncio
from langchain_core.tools import tool
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from config import settings
from tools.image_generator import generate_image
from typing import Any, Union
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# ─── Manus Executive Design System ──────────────────────────────────────────
# Paleta: Dark Navy + Cobalt Blue — limpo, profissional, moderno
_BG     = RGBColor(10,  12,  16)   # Quase preto profundo
_CARD   = RGBColor(28,  33,  40)   # Cinza azulado escuro para elementos
_ACCENT = RGBColor(88, 166, 255)   # Azul cobalto brilhante (Electric Blue)
_WHITE  = RGBColor(255, 255, 255)  # Branco puro
_TEXT   = RGBColor(230, 237, 243)  # Off-white para leitura
_MUTED  = RGBColor(110, 118, 129)  # Cinza médio para detalhes

_W = Inches(13.333)
_H = Inches(7.5)


def _bg(slide, color=_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color, alpha=1.0):
    """Adiciona retângulo colorido."""
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background() # sem borda
    return s


def _add_decorative_elements(slide):
    """Adiciona formas geométricas sutis para visual premium."""
    # Triângulo decorativo no canto inferior direito
    shape = slide.shapes.add_shape(5, _W - Inches(3), _H - Inches(2), Inches(3), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _CARD
    shape.line.fill.background()
    
    # Linha fina de acento
    line = slide.shapes.add_shape(1, Inches(0.5), _H - Inches(0.8), _W - Inches(1), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = _ACCENT


def _text(slide, l, t, w, h, txt, size, bold=False, color=_TEXT,
          align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(txt)
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def _build_cover(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)

    # Grande retângulo central de destaque (levemente mais alto para títulos longos)
    _rect(slide, Inches(0), Inches(2.2), _W, Inches(3.2), _CARD)
    
    # Barra vertical de acento
    _rect(slide, Inches(0.8), Inches(2.2), Pt(8), Inches(3.2), _ACCENT)

    # Título principal - Ajuste dinâmico de fonte para títulos longos
    font_size = 54 if len(title) < 30 else 42
    _text(slide, Inches(1.2), Inches(2.5), Inches(11.0), Inches(1.8),
          title.upper(), font_size, bold=True, color=_WHITE)

    # Subtítulo - Posicionado mais abaixo para não colidir
    if subtitle:
        _text(slide, Inches(1.2), Inches(4.3), Inches(10.5), Inches(0.8),
              subtitle, 22, color=_ACCENT)

    # Footer decorativo
    _text(slide, Inches(0.8), _H - Inches(0.6), Inches(8), Inches(0.4),
          "EXECUTIVE STRATEGY DECK  ·  MANUS AI PRO", 10, color=_MUTED)


def _build_content(prs, title, bullets, img_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _add_decorative_elements(slide)

    # Header moderno
    _rect(slide, Inches(0), Inches(0), _W, Inches(1.1), _CARD)
    _rect(slide, Inches(0.5), Inches(0.35), Pt(4), Inches(0.4), _ACCENT)
    
    _text(slide, Inches(0.75), Inches(0.30), Inches(12), Inches(0.6),
          title.upper(), 28, bold=True, color=_WHITE)

    # Imagem
    has_image = False
    if img_path:
        clean_name = str(img_path)
        tag_match = re.search(r'<SEND_FILE:([^>]+)>', clean_name)
        if tag_match:
            clean_name = tag_match.group(1).strip()
        else:
            clean_name = clean_name.replace("SEND_FILE:", "").replace("<", "").replace(">", "").strip()

        variants = [clean_name, clean_name.replace("-", "_"), clean_name.replace("_", "-")]
        found_path = None
        for v in variants:
            fp = os.path.join(settings.DATA_OUTPUTS_PATH, v)
            if os.path.exists(fp):
                found_path = fp
                break

        if found_path:
            try:
                # Moldura da imagem (Premium Look)
                _rect(slide, Inches(0.45), Inches(1.45), Inches(6.1), Inches(5.1), _ACCENT)
                slide.shapes.add_picture(
                    found_path, Inches(0.5), Inches(1.5),
                    width=Inches(6.0), height=Inches(5.0)
                )
                has_image = True
            except Exception as e:
                logger.error(f"[PPTX] Erro imagem: {e}")

    # Conteúdo (Texto) - Espaçamento de linha melhorado
    txt_l = Inches(7.0) if has_image else Inches(1.0)
    txt_w = Inches(5.8) if has_image else Inches(11.3)
    
    tb = slide.shapes.add_textbox(txt_l, Inches(1.6), txt_w, Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(18) # Mais espaço entre bullets
        p.space_after = Pt(4)
        p.level = 0
        
        # Marcador customizado (Exec Dot)
        r = p.add_run()
        r.text = "◈  "
        r.font.color.rgb = _ACCENT
        r.font.bold = True
        
        r2 = p.add_run()
        r2.text = str(bullet)
        r2.font.size = Pt(19) # Fonte levemente menor para caber melhor
        r2.font.color.rgb = _TEXT


class PpptxSchema(BaseModel):
    slides_content_json: Any = Field(..., description=(
        "Conteúdo dos slides em formato JSON estruturado. "
        "Deve conter 'presentation_title' e uma lista 'slides'. "
        "Cada slide da lista deve ter: 'title' (str) e 'bullets' (lista de str). "
        "DICA PRO: Para criar PowerPoints Premium, gere imagens antes usando o image_generator e insira a "
        "tag (ex: <SEND_FILE:img.png>) dentro de um dos bullets ou no campo 'img_path' do slide."
    ))

@tool(args_schema=PpptxSchema)
async def generate_pptx(slides_content_json: Any) -> str:
    """
    Gera apresentação executiva premium com design Manus AI (Navy + Cobalt Blue, Calibri).
    Aceita string JSON, Dict (ideal) com 'presentation_title' e 'slides', ou Lista de strings.
    """
    if slides_content_json is None:
        logger.error("[PPTXGen] ERRO: slides_content_json é None! Usando dados genéricos.")
        slides_content_json = {"title": "Apresentação", "slides": [{"title": "Aviso", "bullets": ["Conteúdo não fornecido pelo modelo."]}]}

    try:
        filename = f"Exec-Deck-{uuid.uuid4().hex[:6]}.pptx"
        filepath = os.path.join(settings.DATA_OUTPUTS_PATH, filename)

        # ==================================================================
        # NORMALIZAÇÃO DE PARÂMETROS
        # ==================================================================
        prs_title = "EXECUTIVE DECK"
        prs_subtitle = ""
        slides_data = []

        if isinstance(slides_content_json, str):
            try:
                slides_content_json = json.loads(slides_content_json)
            except Exception:
                slides_data = [{"title": "Informação", "bullets": [s.strip() for s in slides_content_json.split('\n') if s.strip()]}]
                
        if isinstance(slides_content_json, list):
            logger.info(f"[PPTXGen] Recebido LISTA com {len(slides_content_json)} itens. Convertendo...")
            prs_title = "Relatório Executivo Automático"
            for item in slides_content_json[:15]:
                if isinstance(item, dict):
                    slides_data.append({
                        "title": str(item.get("title", "Aspecto Identificado")),
                        "bullets": [str(item.get("content", item))]
                    })
                else:
                    slides_data.append({"title": "Ponto Chave", "bullets": [str(item)]})

        elif isinstance(slides_content_json, dict):
            logger.info(f"[PPTXGen] Recebido DICT com chaves: {list(slides_content_json.keys())}")
            prs_title = slides_content_json.get("presentation_title", slides_content_json.get("title", "Apresentação Estratégica"))
            prs_subtitle = slides_content_json.get("subtitle", "")
            
            raw_slides = slides_content_json.get("slides", slides_content_json.get("content", []))
            if isinstance(raw_slides, list):
                for s in raw_slides:
                    if isinstance(s, dict):
                        slides_data.append(s)
                    else:
                        slides_data.append({"title": "Ponto Analítico", "bullets": [str(s)]})
            elif isinstance(raw_slides, str):
                slides_data.append({"title": "Resumo", "bullets": [raw_slides]})
                
        if not slides_data:
            slides_data = [{"title": "Sem Dados Estruturados", "bullets": ["O agente não enviou informações estruturadas adequadas."]}]

        # ==================================================================
        # GERAÇÃO PREMIUM (REMASTERIZADA MANUS AI)
        # ==================================================================
        prs = Presentation()
        prs.slide_width = _W
        prs.slide_height = _H

        # Slide 1: CAPA (Design Navy + Cobalt)
        _build_cover(prs, prs_title, prs_subtitle)
        
        # Slides de Conteúdo
        for i, s_data in enumerate(slides_data):
            title = s_data.get("title", f"PONTO ANALÍTICO {i+1}")
            bullets = s_data.get("bullets", s_data.get("content", []))
            if isinstance(bullets, str): bullets = [bullets]
            
            # --- INTELIGÊNCIA DE IMAGEM ORION ---
            # Escaneia os bullets em busca de tags de imagem geradas pelo Executor
            img_path = s_data.get("image", s_data.get("img_path", None))
            
            if not img_path:
                for b_idx, bullet in enumerate(bullets):
                    match = re.search(r'<SEND_FILE:(img-[^>]+)>', str(bullet))
                    if match:
                        img_path = match.group(1)
                        # Remove a tag do texto para não ficar poluído, já que ela vai pro slide
                        bullets[b_idx] = re.sub(r'<SEND_FILE:[^>]+>', '', str(bullet)).strip()
                        break
            
            # --- AUTO-GERAÇÃO MANUS AI (Fallback) ---
            # Se a IA principal esqueceu de passar a imagem para o PPTX, criamos uma na hora!
            if not img_path and i < 8: # Limitamos a 8 slides automáticos para não travar a API
                try:
                    logger.info(f"[PPTXGen] Slide {i+1}: IA não mandou imagem. Auto-gerando com tema: {title}...")
                    prompt_img = f"Corporate presentation slide background for topic: {title}. Professional, executive, abstract blue data analytics style, minimal."
                    res_img = await generate_image.ainvoke({"prompt": prompt_img, "orientation": "square"})
                    m = re.search(r'<SEND_FILE:(img-[^>]+)>', str(res_img))
                    if m:
                        img_path = m.group(1)
                        logger.info(f"[PPTXGen] Auto-imagem injetada: {img_path}")
                except Exception as e_img:
                    logger.error(f"[PPTXGen] Falhou ao auto-gerar imagem para slide {i+1}: {e_img}")

            _build_content(prs, title, bullets, img_path)

        await asyncio.to_thread(prs.save, filepath)
        
        exists = os.path.exists(filepath)
        size_bytes = os.path.getsize(filepath) if exists else 0
        logger.info(f"[PPTX] Apresentação salva: {filepath} | exists={exists} | size={size_bytes}B")
        return f"Apresentação Executiva gerada: <SEND_FILE:{filename}>"
    except Exception as e:
        logger.error(f"[PPTX] Erro: {e}", exc_info=True)
        return f"Falha no PPTX: {str(e)}"
